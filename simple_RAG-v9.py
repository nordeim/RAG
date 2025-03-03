import os
import sys
import re
import logging
import shutil
import mimetypes
import json
import time
import argparse
from pathlib import Path
from typing import List, Tuple, Optional, Any, Dict, Union
from datetime import datetime

# Core dependencies
import yaml
import requests
import numpy as np
from tqdm import tqdm

# UI
import gradio as gr

# Document processing
try:
    import fitz  # PyMuPDF
except ImportError:
    print("PyMuPDF not found. PDF processing will not be available. Install with 'pip install pymupdf'")

try:
    import pytesseract
    from PIL import Image
except ImportError:
    print("Pytesseract or PIL not found. Image processing will not be available.")

try:
    import aspose.words as aw
except ImportError:
    print("Aspose.Words not found. DOCX processing will be limited.")

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not found. PPTX processing will not be available.")

try:
    import pandas as pd
except ImportError:
    print("Pandas not found. Excel processing will not be available.")

try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
except ImportError:
    print("EbookLib or BeautifulSoup not found. EPUB processing will not be available.")

# LLM and vectorstore
try:
    import ollama
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import Chroma
    from langchain_ollama import OllamaEmbeddings
    from langchain.schema import Document
except ImportError:
    print("LangChain dependencies not found. Install with 'pip install langchain langchain-community langchain-ollama'")
    sys.exit(1)

try:
    from openai import OpenAI
except ImportError:
    print("OpenAI package not found. OpenAI integration will not be available.")

# Configure logging
log_dir = "logs"
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, f"rag_assistant_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log")

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)

logger = logging.getLogger("RAG_Assistant")

# Application constants
APP_NAME = "Enhanced RAG Assistant"
APP_VERSION = "1.0.0"
DATA_DIR = "rag_data"
RAG_SOURCE_FOLDER = os.path.join(DATA_DIR, "sources")
VECTOR_DB_PATH = os.path.join(DATA_DIR, "vector_db")
CONFIG_FILE = os.path.join(DATA_DIR, "config.yaml")
CACHE_DIR = os.path.join(DATA_DIR, "cache")

# Create necessary directories
for directory in [DATA_DIR, RAG_SOURCE_FOLDER, VECTOR_DB_PATH, CACHE_DIR]:
    os.makedirs(directory, exist_ok=True)

# Default configuration
DEFAULT_CONFIG = {
    "ollama": {
        "base_url": "http://localhost:11434",
        "model": "llama2",
        "embedding_model": "nomic-embed-text"
    },
    "openai": {
        "base_url": "https://api.openai.com/v1",
        "api_key": "",
        "model": "gpt-3.5-turbo"
    },
    "ui": {
        "theme": "default",
        "default_provider": "Ollama",
        "max_history": 10
    },
    "rag": {
        "chunk_size": 500,
        "chunk_overlap": 50,
        "k_retrieval": 5,
        "system_prompt": "You are a helpful assistant that provides accurate information based on the given context. If the answer is not in the context, say you don't know."
    },
    "generation": {
        "temperature": 0.7,
        "max_tokens": 1024
    }
}

# Global session state
class SessionState:
    def __init__(self):
        self.vectorstore = None
        self.retriever = None
        self.text_splitter = None
        self.processed_files = set()
        self.config = self.load_config()
        self.query_history = []
        self.response_cache = {}
        self.current_status = "Ready"
        self.last_error = None
    
    def load_config(self) -> Dict:
        """Load configuration from file or create default"""
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                    config = yaml.safe_load(f)
                logger.info(f"Loaded configuration from {CONFIG_FILE}")
                
                # Merge with defaults to ensure all keys exist
                merged_config = DEFAULT_CONFIG.copy()
                for section, values in config.items():
                    if section in merged_config:
                        merged_config[section].update(values)
                    else:
                        merged_config[section] = values
                
                return merged_config
            except Exception as e:
                logger.error(f"Error loading configuration: {str(e)}")
                return DEFAULT_CONFIG.copy()
        else:
            # Save default config
            self.save_config(DEFAULT_CONFIG)
            return DEFAULT_CONFIG.copy()
    
    def save_config(self, config: Dict) -> None:
        """Save configuration to file"""
        try:
            with open(CONFIG_FILE, "w", encoding="utf-8") as f:
                yaml.dump(config, f, default_flow_style=False)
            logger.info(f"Saved configuration to {CONFIG_FILE}")
            self.config = config
        except Exception as e:
            logger.error(f"Error saving configuration: {str(e)}")
    
    def update_status(self, status: str) -> None:
        """Update current status"""
        self.current_status = status
        logger.info(f"Status: {status}")
    
    def reset(self) -> None:
        """Reset session state"""
        self.vectorstore = None
        self.retriever = None
        self.processed_files = set()
        self.query_history = []
        self.response_cache = {}
        self.current_status = "Ready"
        self.last_error = None
        logger.info("Session state reset")
    
    def add_to_history(self, query: str, response: str) -> None:
        """Add query and response to history"""
        max_history = self.config["ui"]["max_history"]
        self.query_history.append({"query": query, "response": response, "timestamp": datetime.now().isoformat()})
        if len(self.query_history) > max_history:
            self.query_history = self.query_history[-max_history:]

# Initialize session state
SESSION = SessionState()

# ---------------------- Utility Functions ----------------------
def get_available_ollama_models() -> List[str]:
    """Get list of available Ollama models"""
    try:
        base_url = SESSION.config["ollama"]["base_url"]
        response = requests.get(f"{base_url}/api/tags")
        if response.status_code == 200:
            models = [model["name"] for model in response.json()["models"]]
            return models
        else:
            logger.error(f"Failed to get Ollama models: {response.text}")
            return ["llama2"]  # Default fallback
    except Exception as e:
        logger.error(f"Error getting Ollama models: {str(e)}")
        return ["llama2"]  # Default fallback

def calculate_md5(file_path: str) -> str:
    """Calculate MD5 hash of a file"""
    import hashlib
    hash_md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()

def format_file_size(size_bytes: int) -> str:
    """Format file size in human-readable format"""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.2f} TB"

def detect_file_type(file_path: str) -> str:
    """Improved MIME type detection with fallback"""
    mime_type, _ = mimetypes.guess_type(file_path)
    if not mime_type:
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            '.epub': 'application/epub+zip',
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.txt': 'text/plain',
            '.csv': 'text/csv',
            '.md': 'text/markdown',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.webp': 'image/webp'
        }
        mime_type = type_map.get(ext, 'application/octet-stream')
    return mime_type

def sanitize_filename(filename: str) -> str:
    """Sanitize filename to remove special characters"""
    return re.sub(r'[\\/*?:"<>|]', "_", filename)

# ---------------------- Document Processing ----------------------
def convert_pdf_to_text(file_path: str) -> str:
    """Extract text from a PDF file."""
    try:
        doc = fitz.open(file_path)
        text = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text.append(page.get_text())
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"Error converting PDF to text: {str(e)}")
        return f"Error extracting text from PDF: {str(e)}"

def convert_pptx_to_text(file_path: str) -> str:
    """Extract text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                text.append(f"Slide {slide_num + 1}: " + "\n".join(slide_text))
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"Error converting PPTX to text: {str(e)}")
        return f"Error extracting text from PPTX: {str(e)}"

def convert_docx_to_text(file_path: str) -> str:
    """Extract text from a DOCX file."""
    try:
        if 'aspose.words' in sys.modules:
            doc = aw.Document(file_path)
            return doc.get_text()
        else:
            # Fallback using python-docx if available
            try:
                import docx
                doc = docx.Document(file_path)
                return "\n\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                return "Aspose.Words and python-docx not available for DOCX conversion."
    except Exception as e:
        logger.error(f"Error converting DOCX to text: {str(e)}")
        return f"Error extracting text from DOCX: {str(e)}"

def convert_xlsx_to_text(file_path: str) -> str:
    """Extract text from an Excel file using pandas."""
    try:
        text_parts = []
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text_parts.append(f"Sheet: {sheet_name}\n{df.to_csv(index=False)}")
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error converting XLSX to text: {str(e)}")
        return f"Error extracting text from XLSX: {str(e)}"

def convert_csv_to_text(file_path: str) -> str:
    """Extract text from a CSV file."""
    try:
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    except Exception as e:
        logger.error(f"Error converting CSV to text: {str(e)}")
        return f"Error extracting text from CSV: {str(e)}"

def convert_epub_to_text(file_path: str) -> str:
    """Extract text from an EPUB file."""
    try:
        book = epub.read_epub(file_path)
        texts = []
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                texts.append(soup.get_text())
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"Error converting EPUB to text: {str(e)}")
        return f"Error extracting text from EPUB: {str(e)}"

def convert_text_file_to_text(file_path: str) -> str:
    """Read plain text from a file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encodings
        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return "Could not decode file with any encoding."
    except Exception as e:
        logger.error(f"Error reading text file: {str(e)}")
        return f"Error reading text file: {str(e)}"

def convert_image_to_text(file_path: str) -> str:
    """Extract text from an image using pytesseract."""
    try:
        img = Image.open(file_path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logger.error(f"Error extracting text from image: {str(e)}")
        return f"Error extracting text from image: {str(e)}"

def convert_to_text(input_file: str) -> str:
    """Convert file to text based on its MIME type."""
    SESSION.update_status(f"Converting {os.path.basename(input_file)} to text...")
    
    converters = {
        'application/pdf': convert_pdf_to_text,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': convert_pptx_to_text,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': convert_docx_to_text,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': convert_xlsx_to_text,
        'application/epub+zip': convert_epub_to_text,
        'text/plain': convert_text_file_to_text,
        'text/csv': convert_csv_to_text,
        'text/markdown': convert_text_file_to_text
    }
    
    mime_type = detect_file_type(input_file)
    
    # Handle image files
    if mime_type and mime_type.startswith('image/'):
        return convert_image_to_text(input_file)
    
    converter = converters.get(mime_type)
    if not converter:
        logger.warning(f"Unsupported file type: {mime_type} for {input_file}")
        return f"Unsupported file type: {mime_type}"
    
    try:
        text = converter(input_file)
        logger.info(f"Successfully converted {input_file}")
        return text
    except Exception as e:
        logger.error(f"Conversion error for {input_file}: {str(e)}")
        return f"Conversion error: {str(e)}"

def extract_metadata(file_path: str) -> Dict[str, Any]:
    """Extract metadata from file"""
    try:
        stat = os.stat(file_path)
        metadata = {
            "filename": os.path.basename(file_path),
            "file_type": detect_file_type(file_path),
            "file_size": stat.st_size,
            "created_time": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "md5": calculate_md5(file_path)
        }
        
        # Extract additional metadata based on file type
        if metadata["file_type"] == "application/pdf":
            try:
                doc = fitz.open(file_path)
                metadata.update({
                    "page_count": len(doc),
                    "title": doc.metadata.get("title", ""),
                    "author": doc.metadata.get("author", ""),
                    "subject": doc.metadata.get("subject", ""),
                    "keywords": doc.metadata.get("keywords", "")
                })
            except Exception as e:
                logger.error(f"Error extracting PDF metadata: {str(e)}")
        
        return metadata
    except Exception as e:
        logger.error(f"Error extracting metadata: {str(e)}")
        return {
            "filename": os.path.basename(file_path),
            "error": str(e)
        }

def process_documents(file_paths: List[str], export_filename: str = None) -> Dict[str, Any]:
    """Process documents and build vector store"""
    if not file_paths:
        return {"status": "error", "message": "No files provided"}
    
    SESSION.update_status("Processing documents...")
    result = {
        "status": "success",
        "processed_files": [],
        "skipped_files": [],
        "errors": []
    }
    
    all_chunks = []
    all_texts = []
    
    # Check for new files
    new_files = [f for f in file_paths if f not in SESSION.processed_files]
    if not new_files:
        SESSION.update_status("No new files to process")
        return {"status": "info", "message": "No new files to process"}
    
    # Initialize text splitter
    chunk_size = SESSION.config["rag"]["chunk_size"]
    chunk_overlap = SESSION.config["rag"]["chunk_overlap"]
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
    )
    SESSION.text_splitter = text_splitter
    
    try:
        # Initialize embeddings
        ollama_base_url = SESSION.config["ollama"]["base_url"]
        ollama_model = SESSION.config["ollama"]["embedding_model"]
        embeddings = OllamaEmbeddings(
            base_url=ollama_base_url,
            model=ollama_model
        )
        
        # Process each file
        for file_path in tqdm(new_files, desc="Processing files"):
            try:
                # Copy file to source folder
                file_name = os.path.basename(file_path)
                sanitized_name = sanitize_filename(file_name)
                saved_path = os.path.join(RAG_SOURCE_FOLDER, sanitized_name)
                shutil.copy2(file_path, saved_path)
                
                # Extract metadata
                metadata = extract_metadata(saved_path)
                
                # Convert to text
                text_content = convert_to_text(saved_path)
                if text_content.startswith("Error") or text_content.startswith("Unsupported"):
                    result["errors"].append({"file": file_name, "error": text_content})
                    continue
                
                # Store for export
                all_texts.append(f"\n\n=== FILE: {file_name} ===\n{text_content}")
                
                # Split into chunks with metadata
                chunks = text_splitter.create_documents(
                    texts=[text_content],
                    metadatas=[{
                        "source": file_name,
                        "file_type": metadata["file_type"],
                        "chunk": i
                    } for i in range(1)]  # This will be expanded by the splitter
                )
                
                # Update chunk metadata
                for i, chunk in enumerate(chunks):
                    chunk.metadata["chunk"] = i + 1
                    chunk.metadata["total_chunks"] = len(chunks)
                
                all_chunks.extend(chunks)
                SESSION.processed_files.add(file_path)
                result["processed_files"].append(file_name)
                
            except Exception as e:
                error_msg = f"Error processing {file_name}: {str(e)}"
                logger.error(error_msg)
                result["errors"].append({"file": file_name, "error": str(e)})
        
        # Export combined text if requested
        if export_filename and all_texts:
            export_path = os.path.join(RAG_SOURCE_FOLDER, sanitize_filename(export_filename))
            with open(export_path, 'w', encoding='utf-8') as f:
                f.write("\n".join(all_texts))
            logger.info(f"Exported processed text to {export_path}")
        
        # Update vector store if we have new chunks
        if all_chunks:
            SESSION.update_status(f"Building vector store with {len(all_chunks)} chunks...")
            
            # Check if we should update existing store or create new one
            if SESSION.vectorstore:
                # Add new documents to existing store
                SESSION.vectorstore.add_documents(all_chunks)
            else:
                # Create new vector store
                vectorstore = Chroma.from_documents(
                    documents=all_chunks,
                    embedding=embeddings,
                    persist_directory=VECTOR_DB_PATH
                )
                SESSION.vectorstore = vectorstore
            
            # Update retriever with new configuration
            k_retrieval = SESSION.config["rag"]["k_retrieval"]
            SESSION.retriever = SESSION.vectorstore.as_retriever(
                search_kwargs={"k": k_retrieval}
            )
            
            SESSION.update_status(f"Vector store updated with {len(all_chunks)} chunks from {len(result['processed_files'])} files")
            result["chunk_count"] = len(all_chunks)
        
    except Exception as e:
        error_msg = f"Document processing failed: {str(e)}"
        logger.error(error_msg)
        SESSION.last_error = error_msg
        return {"status": "error", "message": error_msg}
    
    return result

# ---------------------- LLM Integration ----------------------
def generate_response(
    question: str,
    system_prompt: Optional[str] = None,
    provider: Optional[str] = None
) -> Dict[str, Any]:
    """Generate a response using the configured LLM"""
    
    # Check if we have processed documents
    if not SESSION.vectorstore or not SESSION.retriever:
        return {
            "status": "error", 
            "message": "Please process documents first before asking questions"
        }
    
    # Use defaults from configuration if not specified
    if system_prompt is None:
        system_prompt = SESSION.config["rag"]["system_prompt"]
    
    if provider is None:
        provider = SESSION.config["ui"]["default_provider"]
    
    # Check cache for identical question
    cache_key = f"{question}|{system_prompt}|{provider}"
    if cache_key in SESSION.response_cache:
        logger.info(f"Returning cached response for: {question}")
        return SESSION.response_cache[cache_key]
    
    try:
        SESSION.update_status(f"Retrieving context for: {question}")
        
        # Retrieve relevant documents
        retrieved_docs = SESSION.retriever.invoke(question)
        
        if not retrieved_docs:
            context = "No relevant information found in the documents."
        else:
            # Format context with sources
            context_parts = []
            for i, doc in enumerate(retrieved_docs):
                source = doc.metadata.get("source", "Unknown")
                chunk = doc.metadata.get("chunk", "Unknown")
                total = doc.metadata.get("total_chunks", "Unknown")
                context_parts.append(f"[Document {i+1}: {source} (Chunk {chunk}/{total})]:\n{doc.page_content}")
            
            context = "\n\n".join(context_parts)
        
        SESSION.update_status(f"Generating response using {provider}...")
        
        # Generate response using selected provider
        if provider == "Ollama":
            response = generate_ollama_response(question, context, system_prompt)
        else:  # OpenAI
            response = generate_openai_response(question, context, system_prompt)
        
        # Format the final result
        result = {
            "status": "success",
            "question": question,
            "response": response["content"],
            "model": response["model"],
            "context": context,
            "sources": [doc.metadata.get("source", "Unknown") for doc in retrieved_docs],
            "timestamp": datetime.now().isoformat()
        }
        
        # Cache the result
        SESSION.response_cache[cache_key] = result
        
        # Add to history
        SESSION.add_to_history(question, response["content"])
        
        SESSION.update_status("Ready")
        return result
        
    except Exception as e:
        error_msg = f"Error generating response: {str(e)}"
        logger.error(error_msg)
        SESSION.last_error = error_msg
        return {
            "status": "error",
            "message": error_msg
        }

def generate_ollama_response(question: str, context: str, system_prompt: str) -> Dict[str, str]:
    """Generate response using Ollama API"""
    ollama_base_url = SESSION.config["ollama"]["base_url"]
    ollama_model = SESSION.config["ollama"]["model"]
    temperature = SESSION.config["generation"]["temperature"]
    max_tokens = SESSION.config["generation"]["max_tokens"]
    
    # Format the prompt
    full_prompt = f"Question: {question}\n\nContext:\n{context}\n\nPlease answer the question based on the provided context."
    
    # OpenAI-compatible API endpoint
    endpoint = "/api/chat"
    url = f"{ollama_base_url}{endpoint}"
    
    payload = {
        "model": ollama_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": full_prompt}
        ],
        "options": {
            "temperature": temperature,
            "num_predict": max_tokens
        }
    }
    
    try:
        response = requests.post(url, json=payload, timeout=60)
        response.raise_for_status()
        
        result = response.json()
        return {
            "content": result["message"]["content"],
            "model": ollama_model
        }
    except requests.exceptions.RequestException as e:
        logger.error(f"Ollama API error: {str(e)}")
        if hasattr(e, 'response') and e.response:
            logger.error(f"Response: {e.response.text}")
        raise Exception(f"Ollama API error: {str(e)}")

def generate_openai_response(question: str, context: str, system_prompt: str) -> Dict[str, str]:
    """Generate response using OpenAI API"""
    openai_base_url = SESSION.config["openai"]["base_url"]
    openai_api_key = SESSION.config["openai"]["api_key"]
    openai_model = SESSION.config["openai"]["model"]
    temperature = SESSION.config["generation"]["temperature"]
    max_tokens = SESSION.config["generation"]["max_tokens"]
    
    if not openai_api_key:
        raise Exception("OpenAI API key not configured. Please update your configuration.")
    
    # Format the prompt
    full_prompt = f"Question: {question}\n\nContext:\n{context}\n\nPlease answer the question based only on the provided context."
    
    try:
        client = OpenAI(base_url=openai_base_url, api_key=openai_api_key)
        completion = client.chat.completions.create(
            model=openai_model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": full_prompt}
            ],
            temperature=temperature,
            max_tokens=max_tokens
        )
        
        return {
            "content": completion.choices[0].message.content,
            "model": openai_model
        }
    except Exception as e:
        logger.error(f"OpenAI API error: {str(e)}")
        raise Exception(f"OpenAI API error: {str(e)}")

# ---------------------- UI Functions ----------------------
def clear_session():
    """Reset the application state"""
    SESSION.reset()
    return {
        upload_files: gr.update(value=None),
        upload_status: "Session reset. All data has been cleared."
    }

def update_config(provider, ollama_url, ollama_model, openai_url, openai_key, openai_model, 
                 system_prompt, temperature, max_tokens, chunk_size, chunk_overlap, k_retrieval):
    """Update configuration settings"""
    try:
        # Update configuration
        config = SESSION.config.copy()
        
        # Update provider settings
        config["ui"]["default_provider"] = provider
        
        # Update Ollama settings
        config["ollama"]["base_url"] = ollama_url
        config["ollama"]["model"] = ollama_model
        
        # Update OpenAI settings
        config["openai"]["base_url"] = openai_url
        config["openai"]["api_key"] = openai_key
        config["openai"]["model"] = openai_model
        
        # Update RAG settings
        config["rag"]["system_prompt"] = system_prompt
        config["rag"]["chunk_size"] = int(chunk_size)
        config["rag"]["chunk_overlap"] = int(chunk_overlap)
        config["rag"]["k_retrieval"] = int(k_retrieval)
        
        # Update generation settings
        config["generation"]["temperature"] = float(temperature)
        config["generation"]["max_tokens"] = int(max_tokens)
        
        # Save configuration
        SESSION.save_config(config)
        
        # Check if we need to rebuild the retriever with new params
        if SESSION.vectorstore and (
            int(chunk_size) != SESSION.config["rag"]["chunk_size"] or
            int(chunk_overlap) != SESSION.config["rag"]["chunk_overlap"] or
            int(k_retrieval) != SESSION.config["rag"]["k_retrieval"]
        ):
            # If chunk size/overlap changed, we should rebuild the vector store
            if (int(chunk_size) != SESSION.config["rag"]["chunk_size"] or
                int(chunk_overlap) != SESSION.config["rag"]["chunk_overlap"]):
                return "Configuration updated. You should reprocess your documents with the new chunking parameters."
            
            # If only k_retrieval changed, just update the retriever
            if int(k_retrieval) != SESSION.config["rag"]["k_retrieval"]:
                SESSION.retriever = SESSION.vectorstore.as_retriever(
                    search_kwargs={"k": int(k_retrieval)}
                )
        
        return "Configuration updated successfully."
    except Exception as e:
        logger.error(f"Error updating configuration: {str(e)}")
        return f"Error updating configuration: {str(e)}"

def process_uploaded_files(files):
    """Process uploaded files and update the vector store"""
    if not files:
        return "No files uploaded."
    
    file_paths = [f.name for f in files]
    result = process_documents(file_paths)
    
    if result["status"] == "error":
        return f"Error: {result['message']}"
    
    if result["status"] == "info":
        return result["message"]
    
    processed = len(result["processed_files"])
    errors = len(result["errors"])
    
    if errors > 0:
        error_details = "\n".join([f"- {e['file']}: {e['error']}" for e in result["errors"]])
        return f"Processed {processed} files with {errors} errors.\n\nErrors:\n{error_details}"
    
    return f"Successfully processed {processed} files with {result.get('chunk_count', 0)} chunks."

def ask_question(question):
    """Handle the question asking interaction"""
    if not question.strip():
        return "Please enter a question."
    
    result = generate_response(question)
    
    if result["status"] == "error":
        return f"Error: {result['message']}"
    
    return result["response"]

def update_provider_visibility(provider):
    """Update UI visibility based on selected provider"""
    if provider == "Ollama":
        return gr.update(visible=True), gr.update(visible=False)
    else:
        return gr.update(visible=False), gr.update(visible=True)

def load_query_history():
    """Load query history for display"""
    if not SESSION.query_history:
        return "No query history yet."
    
    history_text = []
    for i, item in enumerate(SESSION.query_history):
        history_text.append(f"Q{i+1}: {item['query']}")
        history_text.append(f"A{i+1}: {item['response']}\n")
    
    return "\n".join(history_text)

def get_file_info():
    """Get information about processed files"""
    if not SESSION.processed_files:
        return "No files processed yet."
    
    info = []
    info.append(f"**Processed Files:** {len(SESSION.processed_files)}")
    
    if SESSION.vectorstore:
        try:
            collection = SESSION.vectorstore._collection
            info.append(f"**Total Chunks:** {collection.count()}")
        except:
            pass
    
    # List files
    info.append("\n**Files:**")
    for i, file_path in enumerate(SESSION.processed_files, 1):
        file_name = os.path.basename(file_path)
        info.append(f"{i}. {file_name}")
    
    return "\n".join(info)

# ---------------------- Gradio UI ----------------------
with gr.Blocks(title=APP_NAME) as interface:
    gr.Markdown(f"# {APP_NAME} v{APP_VERSION}")
    gr.Markdown("An advanced Retrieval Augmented Generation system for document analysis and Q&A")
    
    with gr.Tabs() as tabs:
        # Main Tab - Chat Interface
        with gr.TabItem("Chat"):
            with gr.Row():
                with gr.Column(scale=2):
                    with gr.Group():
                        gr.Markdown("### Upload Documents")
                        upload_files = gr.File(
                            file_count="multiple",
                            label="Upload files",
                            type="filepath"  # Changed from "file" to "filepath"
                        )
                        with gr.Row():
                            process_btn = gr.Button("Process Documents", variant="primary")
                            clear_btn = gr.Button("Clear All", variant="secondary")
                        upload_status = gr.Textbox(label="Status", interactive=False)
                    
                    with gr.Group():
                        gr.Markdown("### Ask Questions")
                        question_input = gr.Textbox(
                            label="Enter your question",
                            placeholder="What information are you looking for?",
                            lines=2
                        )
                        ask_btn = gr.Button("Ask", variant="primary")
                        answer_output = gr.Markdown(label="Answer")
                
                with gr.Column(scale=1):
                    with gr.Group():
                        gr.Markdown("### Document Analysis")
                        doc_info = gr.Markdown("Upload documents to see information")
                        refresh_info_btn = gr.Button("Refresh Info")
                    
                    with gr.Accordion("Query History", open=False):
                        history_display = gr.Markdown("No query history yet.")
                        refresh_history_btn = gr.Button("Refresh History")
                    
                    with gr.Accordion("Status", open=False):
                        status_display = gr.Markdown("System ready.")
        
        # Settings Tab
        with gr.TabItem("Settings"):
            with gr.Group():
                gr.Markdown("### LLM Provider Settings")
                provider_select = gr.Radio(
                    choices=["Ollama", "OpenAI"],
                    label="Select Provider",
                    value=SESSION.config["ui"]["default_provider"]
                )
                
                with gr.Group(visible=(SESSION.config["ui"]["default_provider"] == "Ollama")) as ollama_group:
                    gr.Markdown("#### Ollama Settings")
                    ollama_url = gr.Textbox(
                        label="Ollama API URL",
                        value=SESSION.config["ollama"]["base_url"],
                        placeholder="http://localhost:11434"
                    )
                    # Get available models
                    available_models = get_available_ollama_models()
                    ollama_model = gr.Dropdown(
                        label="Ollama Model",
                        choices=available_models,
                        value=SESSION.config["ollama"]["model"],
                        allow_custom_value=True
                    )
                
                with gr.Group(visible=(SESSION.config["ui"]["default_provider"] == "OpenAI")) as openai_group:
                    gr.Markdown("#### OpenAI Settings")
                    openai_url = gr.Textbox(
                        label="OpenAI API URL",
                        value=SESSION.config["openai"]["base_url"],
                        placeholder="https://api.openai.com/v1"
                    )
                    openai_key = gr.Textbox(
                        label="OpenAI API Key",
                        value=SESSION.config["openai"]["api_key"],
                        placeholder="sk-...",
                        type="password"
                    )
                    openai_model = gr.Textbox(
                        label="OpenAI Model",
                        value=SESSION.config["openai"]["model"],
                        placeholder="gpt-3.5-turbo"
                    )
            
            with gr.Group():
                gr.Markdown("### RAG Settings")
                system_prompt = gr.Textbox(
                    label="System Prompt",
                    value=SESSION.config["rag"]["system_prompt"],
                    lines=3
                )
                with gr.Row():
                    chunk_size = gr.Number(
                        label="Chunk Size",
                        value=SESSION.config["rag"]["chunk_size"],
                        precision=0
                    )
                    chunk_overlap = gr.Number(
                        label="Chunk Overlap",
                        value=SESSION.config["rag"]["chunk_overlap"],
                        precision=0
                    )
                    k_retrieval = gr.Number(
                        label="Number of Retrieved Chunks",
                        value=SESSION.config["rag"]["k_retrieval"],
                        precision=0
                    )
            
            with gr.Group():
                gr.Markdown("### Generation Settings")
                with gr.Row():
                    temperature = gr.Slider(
                        label="Temperature",
                        value=SESSION.config["generation"]["temperature"],
                        minimum=0.0,
                        maximum=2.0,
                        step=0.1
                    )
                    max_tokens = gr.Number(
                        label="Max Tokens",
                        value=SESSION.config["generation"]["max_tokens"],
                        precision=0
                    )
            
            save_config_btn = gr.Button("Save Settings", variant="primary")
            settings_status = gr.Markdown()
        
        # Help Tab
        with gr.TabItem("Help"):
            gr.Markdown("""
            # Enhanced RAG Assistant Help
            
            ## Overview
            
            This application provides a Retrieval Augmented Generation (RAG) system that allows you to:
            
            1. Upload documents in various formats (PDF, DOCX, PPTX, TXT, CSV, EPUB, images)
            2. Process these documents to extract their textual content
            3. Ask questions about the content of these documents
            4. Get AI-generated answers that reference the specific parts of your documents
            
            ## Using the Application
            
            ### 1. Upload and Process Documents
            
            - Click the **Upload Files** button to select files from your computer
            - Click **Process Documents** to extract text and build the knowledge base
            - The status area will show progress and any errors
            
            ### 2. Ask Questions
            
            - Type your question in the input box
            - Click **Ask** to get an answer
            - The system will:
              - Find the most relevant parts of your documents
              - Generate an answer based on this information
              - Show the answer with references to the source documents
            
            ### 3. Configure Settings
            
            In the Settings tab, you can customize:
            
            - **LLM Provider**: Choose between Ollama (local) or OpenAI
            - **RAG Settings**: Adjust how documents are processed and retrieved
            - **Generation Settings**: Control the output generation
            
            ## File Formats
            
            The system supports:
            
            - PDF documents
            - Microsoft Word (DOCX)
            - Microsoft PowerPoint (PPTX)
            - Microsoft Excel (XLSX)
            - Text files (TXT, MD)
            - EPUB e-books
            - CSV data files
            - Images (JPG, PNG) via OCR
            
            ## Troubleshooting
            
            - If document processing fails, check file formats and try again
            - For large documents, processing may take some time
            - If you get unclear answers, try reformulating your question
            - Check the logs for detailed error information
            
            ## Privacy
            
            All processing happens locally or via your configured API endpoints. Your documents are not shared with any third party unless you're using remote API services like OpenAI.
            """)
    
    # Event handlers
    process_btn.click(
        fn=process_uploaded_files,
        inputs=[upload_files],
        outputs=[upload_status]
    )
    
    clear_btn.click(
        fn=clear_session,
        inputs=None,
        outputs=[upload_files, upload_status]
    )
    
    ask_btn.click(
        fn=ask_question,
        inputs=[question_input],
        outputs=[answer_output]
    )
    
    provider_select.change(
        fn=update_provider_visibility,
        inputs=[provider_select],
        outputs=[ollama_group, openai_group]
    )
    
    save_config_btn.click(
        fn=update_config,
        inputs=[
            provider_select,
            ollama_url,
            ollama_model,
            openai_url,
            openai_key,
            openai_model,
            system_prompt,
            temperature,
            max_tokens,
            chunk_size,
            chunk_overlap,
            k_retrieval
        ],
        outputs=[settings_status]
    )
    
    refresh_history_btn.click(
        fn=load_query_history,
        inputs=None,
        outputs=[history_display]
    )
    
    refresh_info_btn.click(
        fn=get_file_info,
        inputs=None,
        outputs=[doc_info]
    )

# ---------------------- Main Application ----------------------
def main():
    """Main application entry point"""
    parser = argparse.ArgumentParser(description=f"{APP_NAME} v{APP_VERSION}")
    parser.add_argument("--host", type=str, default="127.0.0.1", help="Host IP address")
    parser.add_argument("--port", type=int, default=7860, help="Port number")
    parser.add_argument("--debug", action="store_true", help="Enable debug mode")
    parser.add_argument("--share", action="store_true", help="Share the app publicly")
    
    args = parser.parse_args()
    
    logger.info(f"Starting {APP_NAME} v{APP_VERSION}")
    logger.info(f"Host: {args.host}, Port: {args.port}, Debug: {args.debug}, Share: {args.share}")
    
    interface.launch(
        server_name=args.host,
        server_port=args.port,
        debug=args.debug,
        share=args.share
    )

if __name__ == "__main__":
    main()
