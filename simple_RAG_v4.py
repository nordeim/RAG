import os
import sys
import re
import logging
import shutil  # Added for cross-drive file moving
import mimetypes  # Ensure mimetypes is imported
import fitz  # Added for PDF conversion
from typing import List, Tuple, Optional, Any

# Updated imports
import ollama
import gradio as gr
from langchain_community.document_loaders import PyMuPDFLoader  # Unused now; can remove if desired
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma  # Updated import for Chroma
try:
    from langchain_ollama import OllamaEmbeddings  # Updated import for embeddings
except ModuleNotFoundError:
    print("Module 'langchain_ollama' not found. Please install it with 'pip install langchain-ollama'")
    sys.exit(1)
from langchain.schema import Document  # For creating document objects

import pytesseract
from PIL import Image
import aspose.words as aw         # For DOCX conversion (note: license required for commercial use)
from pptx import Presentation       # For PPTX conversion
import pandas as pd                 # For XLSX conversion
from ebooklib import epub           # For EPUB conversion
from bs4 import BeautifulSoup        # For EPUB conversion

# Configure logging
logging.basicConfig(level=logging.INFO)

# Create RAG source folder if missing
RAG_SOURCE_FOLDER = "RAG_source"
os.makedirs(RAG_SOURCE_FOLDER, exist_ok=True)

def detect_file_type(file_path: str) -> Optional[str]:
    """Detect the file type using mimetypes."""
    mime_type, _ = mimetypes.guess_type(file_path)
    if not mime_type:
        _, ext = os.path.splitext(file_path)
        ext_lower = ext.lower()
        if ext_lower in ['.txt', '.log', '.csv']:
            mime_type = 'text/plain'
        elif ext_lower in ['.epub']:
            mime_type = 'application/epub+zip'
        elif ext_lower in ['.pdf']:
            mime_type = 'application/pdf'
        elif ext_lower in ['.docx', '.doc']:
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif ext_lower in ['.pptx', '.ppt']:
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif ext_lower in ['.xlsx', '.xls']:
            mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif ext_lower in ['.jpg', '.jpeg', '.png', '.gif', '.tiff', '.bmp']:
            mime_type = f'image/{ext_lower[1:]}'
    return mime_type

def convert_pptx_to_text(file_path: str) -> str:
    """Convert PPTX files to text."""
    text = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            text.append(f"[SLIDE {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    except Exception as e:
        return f"Error converting PPTX: {e}"
    return "\n".join(text)

def convert_docx_to_text(file_path: str) -> str:
    """Convert DOCX files to text using Aspose.Words."""
    try:
        doc = aw.Document(file_path)
        return doc.get_text()
    except Exception as e:
        return f"Error converting DOCX: {e}"

def convert_xlsx_to_text(file_path: str) -> str:
    """Convert Excel files to text."""
    try:
        dfs = pd.read_excel(file_path, sheet_name=None)
        text = []
        for sheet, df in dfs.items():
            text.append(f"[SHEET: {sheet}]")
            text.append(df.to_string(index=False))
        return "\n\n".join(text)
    except Exception as e:
        return f"Error converting XLSX: {e}"

def convert_epub_to_text(file_path: str) -> str:
    """Convert EPUB files to text."""
    text = []
    try:
        book = epub.read_epub(file_path)
        for item in book.get_items():
            # Use 9 for document items (replacing epub.ITEM_DOCUMENT)
            if item.get_type() == 9:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return "\n\n".join(text)
    except Exception as e:
        return f"Error converting EPUB: {e}"

def convert_image_to_text(file_path: str) -> str:
    """Extract text from images using OCR."""
    try:
        return pytesseract.image_to_string(Image.open(file_path))
    except Exception as e:
        return f"Error converting Image (OCR): {e}"

def convert_text_file_to_text(file_path: str) -> str:
    """Directly read plain text files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Error reading Text file: {e}"

def convert_to_text(input_file: str, output_file: str) -> Optional[str]:
    """Convert various document types to plain text."""
    file_type = detect_file_type(input_file)
    if file_type == "application/pdf":
        try:
            text = ""
            with fitz.open(input_file) as doc:
                for page in doc:
                    text += page.get_text()
        except Exception as e:
            return f"Error converting PDF: {e}"
    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = convert_pptx_to_text(input_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = convert_docx_to_text(input_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        text = convert_xlsx_to_text(input_file)
    elif file_type == "application/epub+zip":
        text = convert_epub_to_text(input_file)
    elif file_type and file_type.startswith("image/"):
        text = convert_image_to_text(input_file)
    elif file_type == "text/plain":
        text = convert_text_file_to_text(input_file)
    else:
        return f"Unsupported file type: {file_type}"
    
    if isinstance(text, str) and text.startswith("Error"):
        print(text)
        return
    
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(text)
        return text
    except Exception as e:
        return f"Error writing to output file: {e}"

def combine_docs(docs: List[Document]) -> str:
    """Combine document contents using newline separator."""
    return "\n\n".join(doc.page_content for doc in docs)

def process_documents(file_paths: List[str]) -> Tuple[Any, Any, Any]:
    """Process multiple documents and store them in the vector database."""
    all_chunks = []
    embeddings = OllamaEmbeddings(model="deepseek-r1:1.5b")
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)

    for file_path in file_paths:
        # Save the original file in the RAG source folder
        file_name = os.path.basename(file_path)
        saved_path = os.path.join(RAG_SOURCE_FOLDER, file_name)
        # Use shutil.move for cross-drive file moving
        shutil.move(file_path, saved_path)

        # Convert the document to text
        temp_text_file = f"temp_{file_name}.txt"
        text_content = convert_to_text(saved_path, temp_text_file)

        # Instead of using PyMuPDFLoader, read the converted text file directly
        try:
            with open(temp_text_file, 'r', encoding='utf-8') as f:
                content = f.read()
            doc = Document(page_content=content)
            chunks = text_splitter.split_documents([doc])
            all_chunks.extend(chunks)
        except Exception as e:
            logging.error("Error reading converted file %s: %s", temp_text_file, e)

    # Store all chunks in the vector database
    vectorstore = Chroma.from_documents(documents=all_chunks, embedding=embeddings, persist_directory="./chroma_db")
    retriever = vectorstore.as_retriever()
    logging.info("Processed %d documents: %d chunks generated", len(file_paths), len(all_chunks))
    return text_splitter, vectorstore, retriever

def export_processed_data(vectorstore: Any, output_file: str):
    """Export processed data from the vector database into a text file."""
    try:
        with open(output_file, "w", encoding="utf-8") as f:
            for doc in vectorstore.get_all_documents():
                f.write(f"{doc.page_content}\n")
        logging.info("Exported processed data to %s", output_file)
    except Exception as e:
        logging.error("Error exporting processed data: %s", e)

def rag_chain(question: str, text_splitter: Any, vectorstore: Any, retriever: Any, system_prompt: str = "") -> str:
    """Generate a response using the RAG pipeline."""
    retrieved_docs = retriever.invoke(question)
    formatted_content = combine_docs(retrieved_docs)
    formatted_prompt = f"System: {system_prompt}\n\nQuestion: {question}\n\nContext: {formatted_content}"
    response = ollama.chat(
        model="deepseek-r1:1.5b",
        messages=[{'role': 'user', 'content': formatted_prompt}]
    )
    return response['message']['content']

def ask_question(file_paths: List[str], question: str, system_prompt: str = "") -> Optional[str]:
    """Handle the end-to-end process of uploading files and answering questions."""
    text_splitter, vectorstore, retriever = process_documents(file_paths)
    result = rag_chain(question, text_splitter, vectorstore, retriever, system_prompt)
    return result

def run_tests():
    """Run basic tests for key functionality."""
    from collections import namedtuple
    Doc = namedtuple("Doc", ["page_content"])
    docs = [Doc("Hello"), Doc("World")]
    combined = combine_docs([Document(page_content=d.page_content) for d in docs])
    assert combined == "Hello\n\nWorld", "combine_docs test failed"
    logging.info("combine_docs test passed")
    print("All tests passed.")

if __name__ == '__main__':
    if len(sys.argv) > 1 and sys.argv[1] == "test":
        run_tests()
    else:
        interface = gr.Interface(
            fn=ask_question,
            inputs=[
                gr.File(label="Upload Files (optional)", file_count="multiple"),
                gr.Textbox(label="Ask a question"),
                gr.Textbox(label="System Prompt (optional)")
            ],
            outputs="text",
            title="Ask Questions About Your Documents",
            description="Use DeepSeek-R1 1.5B to answer your questions about the uploaded documents."
        )
        interface.launch()

