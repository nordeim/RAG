import os
import sys
import re
import logging
import shutil
import mimetypes
import fitz
from typing import List, Tuple, Optional, Any, Dict
import ollama
import gradio as gr
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
try:
    from langchain_ollama import OllamaEmbeddings
except ModuleNotFoundError:
    print("Module 'langchain_ollama' not found. Please install with 'pip install langchain-ollama'")
    sys.exit(1)
from langchain.schema import Document
import pytesseract
from PIL import Image
import aspose.words as aw
from pptx import Presentation
import pandas as pd
from ebooklib import epub
from bs4 import BeautifulSoup
from openai import OpenAI  # Added OpenAI integration
import yaml                    # << new import for configuration handling
import requests                # << new import for Ollama API call

# Configure logging
logging.basicConfig(level=logging.INFO)
RAG_SOURCE_FOLDER = "RAG_source"
os.makedirs(RAG_SOURCE_FOLDER, exist_ok=True)

# Global session state
SESSION_STATE = {
    "vectorstore": None,
    "retriever": None,
    "text_splitter": None,
    "processed_files": set()
}

# Attempt to load OpenAI config if present
openai_base_url = ""
openai_api_key = ""
openai_model = ""
default_system_prompt = "You are a helpful assistant."  # default fallback

config_path = "api_configuration.yaml"
if os.path.exists(config_path):
    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)
    openai_api_key = config.get("API_Key", "")
    openai_base_url = config.get("API_Url", "")
    openai_model = config.get("Model", "")
    # if System_Prompt is provided in the config, use it as default
    default_system_prompt = config.get("System_Prompt", default_system_prompt)

# Default settings for Ollama
default_ollama_base_url = "http://127.0.0.1:11434"
default_ollama_model = "deepseek-r1:1.5b"  # fallback default (user will select from list)

# ---------------------- Core Functions (Improved) ----------------------
def detect_file_type(file_path: str) -> Optional[str]:
    """Improved MIME type detection with fallback"""
    mime_type, _ = mimetypes.guess_type(file_path)
    if not mime_type:
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            '.epub': 'application/epub+zip',
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.txt': 'text/plain',
            '.csv': 'text/plain',
            '.jpg': 'image/jpeg',
            '.png': 'image/png'
        }
        mime_type = type_map.get(ext, 'application/octet-stream')
    return mime_type

def convert_pptx_to_text(file_path: str) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def convert_docx_to_text(file_path: str) -> str:
    """Extract text from a DOCX file using aspose.words."""
    doc = aw.Document(file_path)
    return doc.get_text()

def convert_xlsx_to_text(file_path: str) -> str:
    """Extract text from an Excel file using pandas."""
    try:
        df = pd.read_excel(file_path)
        return df.to_csv(index=False)
    except Exception as e:
        return f"Conversion error: {str(e)}"

def convert_epub_to_text(file_path: str) -> str:
    """Extract text from an EPUB file."""
    book = epub.read_epub(file_path)
    texts = []
    for item in book.get_items():
        if item.get_type() == epub.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            texts.append(soup.get_text())
    return "\n".join(texts)

def convert_text_file_to_text(file_path: str) -> str:
    """Read plain text from a file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def convert_image_to_text(file_path: str) -> str:
    """Extract text from an image using pytesseract."""
    img = Image.open(file_path)
    return pytesseract.image_to_string(img)

def convert_to_text(input_file: str) -> Optional[str]:
    """Improved conversion with error handling"""
    converters = {
        'application/pdf': lambda f: "".join(page.get_text() for page in fitz.open(f)),
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': convert_pptx_to_text,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': convert_docx_to_text,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': convert_xlsx_to_text,
        'application/epub+zip': convert_epub_to_text,
        'text/plain': convert_text_file_to_text
    }
    
    mime_type = detect_file_type(input_file)
    if mime_type.startswith('image/'):
        return convert_image_to_text(input_file)
    
    converter = converters.get(mime_type)
    if not converter:
        return f"Unsupported file type: {mime_type}"
    
    try:
        return converter(input_file)
    except Exception as e:
        return f"Conversion error: {str(e)}"

# ---------------------- New Features Implementation ----------------------
def process_documents(file_paths: List[str], export_filename: str = None) -> None:
    """Enhanced processing with export capability"""
    all_chunks = []
    all_texts = []
    
    # Check if we need to reprocess files
    new_files = [f for f in file_paths if f not in SESSION_STATE["processed_files"]]
    if not new_files:
        return

    try:
        embeddings = OllamaEmbeddings(model="deepseek-r1:1.5b")
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)

        for file_path in new_files:
            # Copy instead of move to handle temp files
            file_name = os.path.basename(file_path)
            saved_path = os.path.join(RAG_SOURCE_FOLDER, file_name)
            shutil.copy(file_path, saved_path)

            # Convert and process
            text_content = convert_to_text(saved_path)
            if text_content.startswith("Conversion error"):
                logging.error(text_content)
                continue

            # Store for export
            all_texts.append(f"\n\n=== FILE: {file_name} ===\n{text_content}")
            
            # Split into chunks
            doc = Document(page_content=text_content)
            chunks = text_splitter.split_documents([doc])
            all_chunks.extend(chunks)
            SESSION_STATE["processed_files"].add(file_path)

        # Handle text export
        if export_filename:
            export_path = os.path.join(RAG_SOURCE_FOLDER, export_filename)
            with open(export_path, 'w', encoding='utf-8') as f:
                f.write("\n".join(all_texts))
            logging.info(f"Exported processed text to {export_path}")

        # Update vector store
        if all_chunks:
            vectorstore = Chroma.from_documents(
                documents=all_chunks,
                embedding=embeddings,
                persist_directory="./chroma_db"
            )
            SESSION_STATE.update({
                "vectorstore": vectorstore,
                "retriever": vectorstore.as_retriever(),
                "text_splitter": text_splitter
            })
    except Exception as e:
        logging.error(f"Processing failed: {str(e)}")

def generate_response(
    question: str,
    system_prompt: str,
    provider: str,
    openai_base_url: str,
    openai_api_key: str,
    openai_model: str,
    ollama_base_url: str,
    ollama_model: str,
    temperature: float,
    max_tokens: int
) -> str:
    """Unified response generator for both providers"""
    if not SESSION_STATE["vectorstore"]:
        return "Error: Please process documents first!"
    
    retrieved_docs = SESSION_STATE["retriever"].invoke(question)
    context = "\n\n".join(doc.page_content for doc in retrieved_docs)
    
    full_prompt = f"System: {system_prompt}\nQuestion: {question}\nContext:\n{context}"
    
    if provider == "Ollama":
        # Use Ollama API call (based on sample_code_for_Ollama_API_call.py)
        endpoint = "/chat/completions"
        url = f"{ollama_base_url}{endpoint}"
        payload = {
            "model": ollama_model,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": full_prompt}
            ],
            "temperature": temperature,
            "max_tokens": max_tokens
        }
        try:
            response = requests.post(url, json=payload)
            if response.status_code == 200:
                result = response.json()
                return result['choices'][0]['message']['content']
            else:
                return f"Error: {response.status_code} {response.text}"
        except Exception as e:
            return f"API Error: {str(e)}"
    else:  # OpenAI call using OpenAI integration
        try:
            client = OpenAI(base_url=openai_base_url, api_key=openai_api_key)
            completion = client.chat.completions.create(
                model=openai_model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": full_prompt}
                ],
                temperature=temperature,
                max_tokens=max_tokens
            )
            return completion.choices[0].message.content
        except Exception as e:
            return f"API Error: {str(e)}"

def clear_session():
    """Reset the application state"""
    SESSION_STATE.update({
        "vectorstore": None,
        "retriever": None,
        "text_splitter": None,
        "processed_files": set()
    })
    return [None]*4  # Clear file inputs

# ---------------------- Gradio Interface ----------------------
with gr.Blocks(title="Enhanced RAG Assistant") as interface:
    gr.Markdown("## 📚 Intelligent Document Analysis with RAG")
    
    with gr.Row():
        with gr.Column(scale=1):
            provider = gr.Radio(
                choices=["Ollama", "OpenAI"],
                label="LLM Provider",
                value="Ollama"
            )
            # Group for OpenAI parameters
            with gr.Group(visible=False) as openai_group:
                openai_base_url_input = gr.Textbox(
                    label="OpenAI API URL",
                    value=openai_base_url,
                    placeholder="https://api.openai.com/v1"
                )
                openai_api_key_input = gr.Textbox(
                    label="OpenAI API Key",
                    type="password",
                    value=openai_api_key,
                    placeholder="Enter OpenAI API key"
                )
                openai_model_input = gr.Textbox(
                    label="OpenAI Model",
                    value=openai_model,
                    placeholder="e.g., gpt-4, gpt-3.5-turbo"
                )
            # Group for Ollama parameters
            with gr.Group(visible=True) as ollama_group:
                ollama_base_url_input = gr.Textbox(
                    label="Ollama Base URL",
                    value=default_ollama_base_url,
                    placeholder="http://127.0.0.1:11434"
                )
                ollama_model_input = gr.Textbox(
                    label="Ollama Model",
                    value=default_ollama_model,
                    placeholder="Select Ollama model (use 'ollama list' to see options)"
                )
            temperature = gr.Slider(
                label="Temperature",
                minimum=0.0,
                maximum=1.0,
                value=0.7,
                step=0.1
            )
            max_tokens = gr.Number(
                label="Max Tokens",
                value=16384,
                precision=0
            )
            export_check = gr.Checkbox(
                label="Export Converted Text",
                info="Save processed text for future use"
            )
            export_name = gr.Textbox(
                label="Export Filename",
                placeholder="processed_text.txt",
                visible=False
            )
            
            # Toggle API parameter groups based on provider selection.
            def toggle_provider(choice):
                return gr.update(visible=(choice=="OpenAI")), gr.update(visible=(choice=="Ollama"))
            
            provider.change(
                toggle_provider,
                inputs=provider,
                outputs=[openai_group, ollama_group]
            )
            export_check.change(
                lambda x: gr.update(visible=x),
                inputs=export_check,
                outputs=export_name
            )

        with gr.Column(scale=2):
            file_input = gr.Files(label="Upload Documents", file_count="multiple")
            system_prompt = gr.Textbox(
                label="System Prompt",
                value=default_system_prompt,
                placeholder="You are a helpful assistant...",
                lines=3
            )
            question = gr.Textbox(
                label="Your Question",
                placeholder="Ask anything about the documents...",
                lines=3
            )
            submit_btn = gr.Button("Ask", variant="primary")
            clear_btn = gr.Button("Clear Session")
            
            output = gr.Textbox(label="AI Response", interactive=False)

    # Event handling
    submit_btn.click(
        fn=process_documents,
        inputs=[file_input, export_name],
        outputs=[]
    ).success(
        fn=generate_response,
        inputs=[
            question,
            system_prompt,
            provider,
            openai_base_url_input,
            openai_api_key_input,
            openai_model_input,
            ollama_base_url_input,
            ollama_model_input,
            temperature,
            max_tokens
        ],
        outputs=output
    )
    
    clear_btn.click(
        fn=clear_session,
        inputs=[],
        outputs=[file_input, question, system_prompt, export_name]
    )

if __name__ == '__main__':
    interface.launch()
